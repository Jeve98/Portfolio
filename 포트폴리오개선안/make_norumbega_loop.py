# -*- coding: utf-8 -*-
"""Generate an editable PPTX core-loop flowchart for Norumbega."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import copy

# Norumbega concept colors (dark-fantasy boss raid: crimson boss glow, charcoal dungeon, parchment text)
CRIMSON = RGBColor(0x9B, 0x23, 0x35)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
DARK = RGBColor(0x1C, 0x1B, 0x22)
GRAY = RGBColor(0x4A, 0x48, 0x4E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_FILL = RGBColor(0xFC, 0xFB, 0xF9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def set_arrowhead(line_elem, end=True, kind="triangle", w="med", l="med"):
    tag = "tailEnd" if not end else "headEnd"
    el = line_elem.find(qn('a:' + tag))
    if el is None:
        el = line_elem.makeelement(qn('a:' + tag), {})
        line_elem.append(el)
    el.set('type', kind)
    el.set('w', w)
    el.set('len', l)

def add_box(x, y, w, h, title, lines, border_color=CRIMSON, fill=LIGHT_FILL, title_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = 0.08
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2.5)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = title
    r0.font.size = Pt(20)
    r0.font.bold = True
    r0.font.color.rgb = title_color or border_color
    r0.font.name = "Malgun Gothic"

    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(13)
        r.font.color.rgb = DARK
        r.font.name = "Malgun Gothic"

    return shape

def add_connector(x1, y1, x2, y2, color=CRIMSON, width=2.25, dashed=False, label=None, label_pos=None, curve=False):
    conn_type = MSO_CONNECTOR.CURVE if curve else MSO_CONNECTOR.STRAIGHT
    conn = slide.shapes.add_connector(conn_type, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if dashed:
        ln = conn.line._get_or_add_ln()
        d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(d)
    ln = conn.line._get_or_add_ln()
    set_arrowhead(ln, end=True, kind="triangle", w="med", l="med")

    if label:
        lx, ly = label_pos if label_pos else ((x1 + x2) / 2, (y1 + y2) / 2)
        tb = slide.shapes.add_textbox(Inches(lx - 0.6), Inches(ly - 0.18), Inches(1.2), Inches(0.36))
        tf = tb.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(12.5)
        r.font.bold = True
        r.font.color.rgb = color
        r.font.name = "Malgun Gothic"
    return conn

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.6))
tf = title_box.text_frame
p = tf.paragraphs[0]
r = p.add_run()
r.text = "노룸베가 — 코어 루프"
r.font.size = Pt(26)
r.font.bold = True
r.font.color.rgb = DARK
r.font.name = "Malgun Gothic"

# Box geometry (2x2 cycle, clockwise)
BW, BH = 4.0, 2.15
TOP_Y = 1.35
BOT_Y = 4.65
LEFT_X = 0.9
RIGHT_X = 8.45

box_dunning = add_box(
    LEFT_X, TOP_Y, BW, BH,
    "던닝  (일일 5회)",
    ["방어구 / 장신구 드랍", "패시브 스탯 · 속성 옵션 확인"],
    border_color=GOLD,
)

box_boss = add_box(
    RIGHT_X, TOP_Y, BW, BH,
    "보스레이드",
    ["보스 선택 → 기믹 파악 → 클리어", "무기(액티브 스킬) + 결정 드랍"],
    border_color=CRIMSON,
)

box_build = add_box(
    RIGHT_X, BOT_Y, BW, BH,
    "빌드 강화",
    ["결정 → 무기 스킬 재롤", "토큰 → 무기 수치 재롤"],
    border_color=CRIMSON,
)

box_higher = add_box(
    LEFT_X, BOT_Y, BW, BH,
    "상위 난이도 도전",
    ["더 높은 난이도의 보스", "→ 더 높은 등급의 아이템"],
    border_color=GOLD,
)

# Primary loop arrows (clockwise)
# 던닝 -> 보스레이드 (top edge)
add_connector(
    LEFT_X + BW, TOP_Y + BH / 2,
    RIGHT_X, TOP_Y + BH / 2,
    color=CRIMSON, width=2.5,
    label="보완", label_pos=(LEFT_X + BW + (RIGHT_X - LEFT_X - BW) / 2, TOP_Y + BH / 2 - 0.32),
)

# 보스레이드 -> 빌드 강화 (right edge)
add_connector(
    RIGHT_X + BW / 2, TOP_Y + BH,
    RIGHT_X + BW / 2, BOT_Y,
    color=CRIMSON, width=2.5,
)

# 빌드 강화 -> 상위 난이도 도전 (bottom edge)
add_connector(
    RIGHT_X, BOT_Y + BH / 2,
    LEFT_X + BW, BOT_Y + BH / 2,
    color=CRIMSON, width=2.5,
)

# 상위 난이도 도전 -> 던닝 (left edge, closes the loop)
add_connector(
    LEFT_X + BW / 2, BOT_Y,
    LEFT_X + BW / 2, TOP_Y + BH,
    color=CRIMSON, width=2.5,
    label="반복", label_pos=(LEFT_X + BW / 2 + 0.55, TOP_Y + BH + (BOT_Y - TOP_Y - BH) / 2),
)

# Secondary resource path: 던닝의 토큰이 빌드 강화로 직접 공급 (dashed diagonal)
add_connector(
    LEFT_X + BW * 0.75, TOP_Y + BH,
    RIGHT_X + BW * 0.25, BOT_Y,
    color=GOLD, width=2.0, dashed=True,
    label="토큰", label_pos=((LEFT_X + BW * 0.75 + RIGHT_X + BW * 0.25) / 2 - 0.3, (TOP_Y + BH + BOT_Y) / 2),
)

# Legend
legend_y = 6.95
leg1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(legend_y), Inches(0.35), Inches(0.08))
leg1.fill.solid(); leg1.fill.fore_color.rgb = CRIMSON; leg1.line.fill.background()
leg1_t = slide.shapes.add_textbox(Inches(1.35), Inches(legend_y - 0.14), Inches(2.2), Inches(0.35))
leg1_t.text_frame.paragraphs[0].add_run().text = "메인 루프 흐름"
leg1_t.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
leg1_t.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
leg1_t.text_frame.paragraphs[0].runs[0].font.name = "Malgun Gothic"

leg2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.6), Inches(legend_y), Inches(0.35), Inches(0.08))
leg2.fill.solid(); leg2.fill.fore_color.rgb = GOLD; leg2.line.fill.background()
leg2_t = slide.shapes.add_textbox(Inches(4.05), Inches(legend_y - 0.14), Inches(2.6), Inches(0.35))
leg2_t.text_frame.paragraphs[0].add_run().text = "보조 자원 공급 (토큰)"
leg2_t.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
leg2_t.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK
leg2_t.text_frame.paragraphs[0].runs[0].font.name = "Malgun Gothic"

prs.save("포트폴리오개선안/norumbega_core_loop.pptx")
print("saved")
