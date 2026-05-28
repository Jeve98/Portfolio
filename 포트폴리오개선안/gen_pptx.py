from pptx import Presentation
from pptx.util import Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import lxml.etree as etree

# ── 팔레트 ─────────────────────────────────────────────────────────────────────
BG      = RGBColor(0x1C, 0x1C, 0x2E)
CARD    = RGBColor(0x28, 0x28, 0x40)
ACCENT  = RGBColor(0x4E, 0xCD, 0xC4)
YELLOW  = RGBColor(0xFF, 0xD9, 0x3D)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0xAA, 0xAA, 0xBB)
RED     = RGBColor(0xFF, 0x6B, 0x6B)
GREEN   = RGBColor(0x6B, 0xD9, 0x8A)
SEC_BG  = RGBColor(0x10, 0x10, 0x20)

KR   = "맑은 고딕"
CODE = "Consolas"

prs = Presentation()
prs.slide_width  = Cm(33.867)
prs.slide_height = Cm(19.05)
BLANK = prs.slide_layouts[6]

# ── 유틸 ────────────────────────────────────────────────────────────────────────
def new_slide():
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return s

def box(slide, text, x, y, w, h,
        size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
        bg=None, italic=False, font=KR, wrap=True):
    txb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    if bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name   = font
    run.font.size   = Pt(size)
    run.font.color.rgb = color
    run.font.bold   = bold
    run.font.italic = italic
    return txb

def rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def multiline(slide, lines, x, y, w, h,
              size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
              bg=None, font=KR, line_spacing=None):
    txb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    if bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for text, cfg in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = cfg.get("align", align)
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
        run = p.add_run()
        run.text = text
        run.font.name  = cfg.get("font", font)
        run.font.size  = Pt(cfg.get("size", size))
        run.font.color.rgb = cfg.get("color", color)
        run.font.bold  = cfg.get("bold", bold)
        run.font.italic = cfg.get("italic", False)
    return txb

def std_title(slide, text, sub=None):
    rect(slide, 0, 0, 33.867, 0.18, ACCENT)
    box(slide, text, 0.8, 0.4, 30, 1.5, size=26, bold=True, color=WHITE)
    if sub:
        box(slide, sub, 0.8, 1.8, 30, 0.9, size=14, color=ACCENT)

def section_slide(title, subtitle, accent_color=ACCENT):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = SEC_BG
    rect(s, 0, 7.5, 33.867, 4, accent_color)
    box(s, title,    0.8, 7.8,  32, 2.5, size=44, bold=True, color=WHITE,   align=PP_ALIGN.CENTER)
    box(s, subtitle, 0.8, 10.5, 32, 1.2, size=18, color=SEC_BG, align=PP_ALIGN.CENTER)
    return s

def card(slide, title, body_lines, x, y, w, h, title_color=ACCENT):
    rect(slide, x, y, w, h, CARD)
    rect(slide, x, y, w, 0.06, title_color)
    box(slide, title, x+0.3, y+0.15, w-0.6, 0.8, size=13, bold=True, color=title_color)
    body = "\n".join(body_lines)
    box(slide, body, x+0.3, y+0.9, w-0.6, h-1.1, size=11, color=WHITE)

def metric_pair(slide, label, before, after, x, y, w=8):
    box(slide, label, x, y,      w, 0.8, size=12, color=GRAY,   bold=False)
    box(slide, before, x, y+0.7, w, 1.5, size=28, color=RED,    bold=True, align=PP_ALIGN.CENTER)
    box(slide, "↓", x, y+2.1,   w, 0.9, size=20, color=GRAY,   bold=False, align=PP_ALIGN.CENTER)
    box(slide, after,  x, y+2.9, w, 1.5, size=28, color=GREEN,  bold=True, align=PP_ALIGN.CENTER)

def role_table(slide, rows, x=0.8, y=3.0):
    headers = ["역할", "비중", "내용"]
    col_w   = [5.0, 3.5, 21.5]
    row_h   = 1.05
    hdr_bg  = ACCENT
    hdr_fg  = BG
    hy = y
    cx = x
    for i, (h, w) in enumerate(zip(headers, col_w)):
        rect(slide, cx, hy, w, 0.75, hdr_bg)
        box(slide, h, cx+0.15, hy+0.1, w-0.3, 0.6, size=12, color=hdr_fg, bold=True, align=PP_ALIGN.CENTER)
        cx += w
    for ri, row in enumerate(rows):
        ry = hy + 0.75 + ri * row_h
        bg_row = RGBColor(0x25, 0x25, 0x38) if ri % 2 == 0 else CARD
        cx = x
        for ci, (cell, w) in enumerate(zip(row, col_w)):
            rect(slide, cx, ry, w, row_h-0.05, bg_row)
            align = PP_ALIGN.CENTER if ci < 2 else PP_ALIGN.LEFT
            fc = ACCENT if ci == 0 else WHITE
            box(slide, cell, cx+0.15, ry+0.12, w-0.3, row_h-0.2,
                size=11, color=fc, align=align)
            cx += w

# ══════════════════════════════════════════════════════════════════════════════
# 1. 표지
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 33.867, 19.05, BG)
rect(s, 0, 0, 0.5, 19.05, ACCENT)
rect(s, 0, 8.5, 33.867, 0.06, ACCENT)

box(s, "전연욱", 2.0, 3.5, 30, 3.5, size=72, bold=True, color=WHITE)
box(s, "게임 클라이언트 프로그래머",
    2.0, 7.2, 30, 1.3, size=24, color=ACCENT, bold=True)
box(s, "Unity / C#  ·  Phaser  ·  Django  ·  Vue.js",
    2.0, 9.0, 30, 0.9, size=16, color=GRAY)
box(s, "삼성 청년 SW 아카데미(SSAFY) 수료",
    2.0, 10.0, 30, 0.9, size=16, color=GRAY)
box(s, "jadam41817@gmail.com",
    2.0, 11.2, 30, 0.9, size=14, color=GRAY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# 2. 프로젝트 개요
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
std_title(s, "개발 프로젝트", "3개 프로젝트 — Unity · Phaser · Django")

projects = [
    ("헤이터",
     "2026-04-22 ~ 05-21  |  6인  |  Unity, Spring, React",
     ["전략 디펜스 웹 게임", "기획 총괄 + 6개 파트 구현", "WebGL 최적화 포함"]),
    ("마네킹: 헌티드 몰",
     "2026-02-20 ~ 03-25  |  6인  |  Unity, Photon PUN2",
     ["비대칭 멀티플레이 PVP", "살인마 진영 + 렌더링·네트워크 최적화", "Photon 멀티플레이 실전"]),
    ("프라가라흐",
     "2025-09-30 ~ 12-25  |  2인  |  Phaser, Django, Vue",
     ["2D 웹 액션 RPG", "클라이언트 + 서버 담당", "AI 코드 구조화 경험"]),
]
px = 0.8
for title, period, bullets in projects:
    rect(s, px, 3.0, 10.5, 14.5, CARD)
    rect(s, px, 3.0, 10.5, 0.08, ACCENT)
    box(s, title,  px+0.3, 3.2,  9.9, 1.0, size=18, bold=True, color=ACCENT)
    box(s, period, px+0.3, 4.3,  9.9, 0.7, size=10, color=GRAY)
    by = 5.2
    for b in bullets:
        box(s, "•  " + b, px+0.3, by, 9.9, 0.8, size=12, color=WHITE)
        by += 0.85
    px += 10.9

# ══════════════════════════════════════════════════════════════════════════════
# 헤이터 섹션
# ══════════════════════════════════════════════════════════════════════════════
section_slide("헤이터", "전략 디펜스 웹 게임  |  Unity C# · Spring · React  |  6인  |  4주", ACCENT)

# 3. 헤이터 — 담당 역할
s = new_slide()
std_title(s, "헤이터 — 담당 역할")
role_table(s, [
    ["아키텍처 설계",  "본인", "Thread safe Singleton 제네릭 베이스 클래스 (CRTP 패턴)"],
    ["세션 관리",      "본인", "새로고침 대응 세션 복구 시스템 (OCP 리팩토링)"],
    ["디펜스 기능",    "본인", "유닛 배치 동기화 · 사정거리 인디케이터 · 오브젝트 풀 스폰"],
    ["탐사 컨텐츠",   "본인 (N종 중 4종)", "기근 · 마력 안개 · 인사 이동 · 저주 정화 이벤트 구현"],
    ["보안",          "본인", "AES-GCM + RSA 세션키 교환 기반 암호화 통신"],
    ["최적화",        "본인", "GPU Instancing · SpriteSync · 픽셀아트 설정으로 WebGL 개선"],
])

# 4. 헤이터 — Thread safe Singleton
s = new_slide()
std_title(s, "헤이터 — Thread safe Singleton", "CRTP 패턴 + 상호 배제")
card(s, "문제",
     ["매니저 클래스마다 싱글톤 생성 코드를 반복 작성 → 유지 보수성 저하"],
     0.8, 3.0, 15.5, 2.5)
card(s, "해결",
     ["제네릭 베이스 클래스로 싱글톤 로직 분리",
      "CRTP 패턴으로 매니저마다 독립 정적 인스턴스 → 캐스팅 없이 직접 접근",
      "lock 블록으로 상호 배제 적용 (멀티 스레드 확장 대비)"],
     0.8, 6.0, 15.5, 4.0)
card(s, "추가 트러블슈팅",
     ["부모 오브젝트로 묶었을 때 DontDestroyOnLoad 동작 안 함",
      "Find + 캐싱 방식은 싱글톤을 클라이언트 코드에서 재구현하는 것과 동일 → 안전성 저하",
      "해결: 생성 시 부모 관계를 끊는 방식 채택"],
     0.8, 10.5, 15.5, 4.0)
code = (
    "public class Singleton<T> : MonoBehaviour\n"
    "    where T : MonoBehaviour\n"
    "{\n"
    "    private static T _instance;\n"
    "    private static readonly object _lock\n"
    "                     = new object();\n"
    "    public static T Instance\n"
    "    {\n"
    "        get { lock (_lock) { /* 생성 */ } }\n"
    "    }\n"
    "}"
)
box(s, code, 17.0, 3.0, 16.0, 11.5, size=12, font=CODE, color=ACCENT, bg=CARD)

# 5. 헤이터 — 세션 복구
s = new_slide()
std_title(s, "헤이터 — 세션 복구 시스템", "OCP 위반 → Dictionary 기반 리팩토링")
card(s, "문제",
     ["웹 게임 새로고침 시 씬 재로드",
      "상태 미저장 → 다음 스테이지 정보 미리 확인 / 방어전 무한 반복 가능"],
     0.8, 3.0, 15.5, 2.8)
multiline(s, [
    ("Before  —  switch문 (OCP 위반)", {"size": 12, "color": ACCENT, "bold": True}),
    ("새 상태 추가 시 switch 수정 필요 → 확장 폐쇄 원칙 위반", {"size": 11, "color": GRAY}),
    ("After  —  Dictionary 상태-핸들러 매핑", {"size": 12, "color": GREEN, "bold": True}),
    ("딕셔너리에 항목만 추가 → 기존 코드 변경 없음", {"size": 11, "color": GRAY}),
], 0.8, 6.5, 15.5, 5.0)
code = (
    "void Awake()\n"
    "{\n"
    "    Action recover =\n"
    "      () => StartCoroutine(\n"
    "        RecoveryTransitionRoutine(\"Defense\"));\n\n"
    "    _handlers = new Dictionary<GameState, Action>\n"
    "    {\n"
    "        { GameState.Preparation, recover },\n"
    "        { GameState.Battle,      recover },\n"
    "        { GameState.PostStage,   LoadPostStage },\n"
    "    };\n"
    "}\n\n"
    "public void CheckAndRecoverySession()\n"
    "{\n"
    "    var state = GetLastSavedState();\n"
    "    if (_handlers.TryGetValue(state,\n"
    "           out Action handler))\n"
    "        handler();\n"
    "}"
)
box(s, code, 17.0, 3.0, 16.0, 14.5, size=11, font=CODE, color=ACCENT, bg=CARD)

# 6. 헤이터 — 유닛 배치 + 사정거리
s = new_slide()
std_title(s, "헤이터 — 유닛 배치 & 사정거리 인디케이터")
card(s, "유닛 배치 — 단방향 동기화 + 상태 보존",
     ["문제: UI 변경마다 오브젝트 생성/파괴 → 버프·HP 등 컴포넌트 상태 초기화",
      "해결: UIGridController에서 배치 상태 저장,",
      "      패널 닫힐 때 SyncDeploymentWithWorld() 호출",
      "      고유 ID 매핑으로 재배치 시 이동(파괴 X) → 상태 보존"],
     0.8, 3.0, 15.5, 6.5)
card(s, "사정거리 인디케이터 — 상대적 비율 계산",
     ["문제: 유닛마다 고정값 저장 시 서버 밸런스 변경 때 클라이언트도 수정 필요",
      "해결: 가장 짧은 사정거리(방패병) 기준으로",
      "      인디케이터 크기를 상대 비율로 계산 → 서버 독립적",
      "버프 증가분은 별도 레이어로 추가 표시"],
     0.8, 10.0, 15.5, 5.5)
code = (
    "// 사정거리 인디케이터\n"
    "float baseScale =\n"
    "    (unitRange / shieldRange) * 1.4f;\n"
    "float effectiveRange =\n"
    "    unitRange * (1f + percentSum) + flatSum;\n"
    "float effectiveScale =\n"
    "    (effectiveRange / shieldRange) * 1.4f;\n\n"
    "if (effectiveScale > baseScale)\n"
    "    SpawnIndicatorObject(\n"
    "        bonusIndicator, coord, effectiveScale);\n"
    "SpawnIndicatorObject(\n"
    "    baseIndicator, coord, baseScale);"
)
box(s, code, 17.0, 3.0, 16.0, 10.0, size=12, font=CODE, color=ACCENT, bg=CARD)

# 7. 헤이터 — 오브젝트 풀
s = new_slide()
std_title(s, "헤이터 — 오브젝트 풀 기반 스폰 시스템", "Pending Queue 방식")
card(s, "배경",
     ["적 유닛 300 ~ 1,500+ 마리 등장",
      "웹 게임: 브라우저 리소스 공유 → GC 비용에 민감",
      "유닛 생성/파괴 반복 시 GC 스파이크 빈번"],
     0.8, 3.0, 14.0, 4.5)
card(s, "설계",
     ["웨이브 등장 유닛의 40% 사전 풀 생성",
      "풀 고갈 시 → Pending Queue에 적재",
      "유닛 반납 이벤트 발생 시 → Pending 처리"],
     0.8, 8.0, 14.0, 5.5)
code = (
    "private void SpawnEnemy(int unitId, int grade)\n"
    "{\n"
    "    var obj = pool.Get(unitId);\n"
    "    if (obj == null)\n"
    "    {\n"
    "        _pendingSpawn[unitId]++;  // pending\n"
    "        return;\n"
    "    }\n"
    "    ApplyGrade(obj, unitId, grade);\n"
    "}\n\n"
    "// 반납 이벤트 시 pending 처리\n"
    "private void OnEnemyReturnedHandler(int id)\n"
    "{\n"
    "    if (_pendingSpawn[id] <= 0) return;\n"
    "    _pendingSpawn[id]--;\n"
    "    SpawnEnemy(id, _unitGrades[id]);\n"
    "}"
)
box(s, code, 15.5, 3.0, 17.5, 13.0, size=12, font=CODE, color=ACCENT, bg=CARD)

# 8. 헤이터 — 탐사 컨텐츠
s = new_slide()
std_title(s, "헤이터 — 탐사 컨텐츠 구현 (4종)", "다회차 랜덤성 부여 선택지 기반 컨텐츠")
events = [
    ("기근",     "지속 디버프",
     ["n회 방어전 동안 소비 식량 2배",
      "횟수 차감을 식량 소비(방어전 시작)와 동기화",
      "→ 씬 재로드 시 횟수/식량 불일치 버그 방지"]),
    ("마력 안개", "연출 + 스탯 분기",
     ["성공/실패 → 버프/디버프 적용",
      "결과 무관 → 다음 전투에 안개 연출 등장",
      "채널 분리로 3가지 분기 처리 단순화",
      "WebGL 쉐이더 미지원 → Unlit 기반 머티리얼 제작"]),
    ("인사 이동", "병과 교체",
     ["보유 병과 중 랜덤 선택 → 다른 병과로 변경",
      "클램프 처리로 보유 수량 초과 방지",
      "부족 병과 보충 또는 복불복 역할"]),
    ("저주 정화", "엘리트 유닛 소환",
     ["실패 시 다음 방어전에 엘리트 추가 등장",
      "로그 분석: 고정 유닛 지정 시 웨이브 조합에 따라 위험도 편차 과대",
      "→ 다음 웨이브 구성에서 무작위 추출 후 엘리트 대체"]),
]
ex = 0.8
for title, tag, bullets in events:
    rect(s, ex, 3.0, 7.9, 13.5, CARD)
    rect(s, ex, 3.0, 7.9, 0.07, ACCENT)
    box(s, title, ex+0.2, 3.15, 5.0, 0.85, size=14, bold=True, color=ACCENT)
    box(s, tag,   ex+5.2, 3.15, 2.5, 0.85, size=10, color=YELLOW, align=PP_ALIGN.RIGHT)
    by = 4.3
    for b in bullets:
        box(s, "•  " + b, ex+0.2, by, 7.3, 0.85, size=11, color=WHITE)
        by += 0.9
    ex += 8.25

# 9. 헤이터 — 암호화 통신
s = new_slide()
std_title(s, "헤이터 — 암호화 통신", "AES-GCM + RSA 세션키 교환")
card(s, "배경",
     ["웹 게임: 개발자 도구로 통신 내역 노출",
      "주요 데이터를 API로 받아 덮어쓰는 치팅 방지 구조 사용 중",
      "→ 통신 보안이 핵심"],
     0.8, 3.0, 32.0, 2.8)
options = [
    ("평문 전송",              "—",    "높음",  "✗", RED),
    ("비대칭키 단독 (RSA)",    "느림", "낮음",  "✗", RED),
    ("대칭키 단독 (AES)",      "빠름", "높음",  "✗", RED),
    ("AES-GCM + RSA 세션키 교환", "빠름", "낮음", "✓", GREEN),
]
headers = ["방식", "속도", "키 노출", "채택"]
col_w   = [14.0, 5.0, 5.0, 3.0]
hx, hy  = 0.8, 6.5
for i, (h, w) in enumerate(zip(headers, col_w)):
    rect(s, hx, hy, w, 0.7, ACCENT)
    box(s, h, hx+0.15, hy+0.1, w-0.3, 0.55, size=12, bold=True, color=BG, align=PP_ALIGN.CENTER)
    hx += w
for ri, (label, speed, risk, adopt, ac) in enumerate(options):
    ry  = hy + 0.7 + ri * 1.1
    rbg = RGBColor(0x30, 0x30, 0x48) if ri % 2 == 0 else CARD
    hx  = 0.8
    vals = [label, speed, risk, adopt]
    for vi, (v, w) in enumerate(zip(vals, col_w)):
        rect(s, hx, ry, w, 1.05, rbg)
        fc = ac if vi == 3 else (YELLOW if ri == 3 else WHITE)
        box(s, v, hx+0.15, ry+0.2, w-0.3, 0.7, size=12, color=fc, align=PP_ALIGN.CENTER)
        hx += w
box(s, "• 게임 실행 시 신규 세션키 생성 → 서버 공개키(RSA)로 암호화하여 교환",
    0.8, 11.5, 32.0, 0.75, size=12, color=WHITE)
box(s, "• 이후 모든 통신은 AES-GCM으로 처리 / IV를 매 요청마다 갱신 → 재전송 공격 방지",
    0.8, 12.3, 32.0, 0.75, size=12, color=WHITE)

# 10. 헤이터 — WebGL 최적화
s = new_slide()
std_title(s, "헤이터 — WebGL 렌더링 최적화")
metric_pair(s, "드로우콜",    "3,500", "180",  1.0,  3.0)
metric_pair(s, "메모리 사용", "60 mb", "15 mb", 10.0, 3.0)
card(s, "① GPU Instancing",
     ["환경 오브젝트 머티리얼 4개 활성화",
      "동일 메시·머티리얼 → 1회 인스턴싱 호출",
      "+ 텍스처·노멀맵 크기 축소 / shadowmask 거리 감소"],
     19.5, 3.0, 13.5, 5.0)
card(s, "② SpriteSync",
     ["유닛: 마스터 1개 + 자식 12개 레이어",
      "기존: 13개 각각에 Animator → CPU 과부하",
      "채택: 마스터에만 Animator, 자식은",
      "      LateUpdate에서 sprite·flipX·color 복사",
      "→ 기존 애니메이션 수정 없이 적용"],
     19.5, 8.5, 13.5, 6.5)
card(s, "③ 픽셀아트 렌더링",
     ["Filter Mode: Bilinear → Point",
      "압축: 무손실 적용",
      "PPU 전체 통일 → 픽셀 흐림 + 충돌 크기 불일치 동시 해결"],
     19.5, 15.5, 13.5, 3.0)

# ══════════════════════════════════════════════════════════════════════════════
# 마네킹 섹션
# ══════════════════════════════════════════════════════════════════════════════
section_slide("마네킹: 헌티드 몰",
              "비대칭 멀티플레이 PVP  |  Unity C# · Photon PUN2  |  6인  |  5주",
              RGBColor(0xFF, 0x6B, 0x6B))

# 11. 마네킹 — 역할
s = new_slide()
std_title(s, "마네킹 — 담당 역할")
role_table(s, [
    ["기획 참여",        "본인",  "살인마 스킬 기획 및 밸런스 논의"],
    ["살인마 진영 구현", "본인",  "더미 설치 · 전이 · 교란 · 위장 스킬 4종"],
    ["스폰 시스템 최적화","본인", "스폰 데이터 구조 재설계 → 네트워크 트래픽 50% 감소"],
    ["상호작용 시스템",  "본인",  "IInteractable 인터페이스 기반 설계"],
    ["증강 시스템",      "본인",  "AbilityBase 추상 클래스 기반 설계"],
    ["렌더링 최적화",    "본인",  "Occlusion Culling 적용 (13fps → 63.6fps, 약 390% 향상)"],
    ["네트워크 안정화",  "본인",  "Jitter 원인 분석 및 동기화 기준 일원화"],
])

# 12. 마네킹 — 기획 의도
s = new_slide()
std_title(s, "마네킹 — 기획 의도")
card(s, "트렌드 분석",
     ["게이머 평균 플레이 타임 지속 감소",
      "롤 칼바람 나락에 증강 추가 → 단판 템포 가속 트렌드",
      "비대칭 PVP(협력·심리전)에 로그라이크 요소(증강) 결합"],
     0.8, 3.0, 15.5, 5.5)
card(s, "AI 활용 팀 규칙",
     ["단순 '만들어줘'가 아닌 명확한 근거와 동작 방식 명시",
      "기존 코드와의 정합성·연동성을 고려하여 프롬프팅",
      "생성된 코드 반드시 직접 검토 후 사용"],
     0.8, 9.2, 15.5, 5.0)
card(s, "증강 시스템",
     ["매 판마다 증강 조합이 달라 → 리플레이어빌리티 확보",
      "로그라이크 장르의 핵심 '빌드 다양성'을 PVP에 이식"],
     17.5, 3.0, 15.5, 5.5)
card(s, "결과",
     ["단판 플레이 템포 가속",
      "매 판 새로운 빌드 경험으로 반복 플레이 유도"],
     17.5, 9.2, 15.5, 5.0)

# 13. 마네킹 — 살인마 스킬 4종
s = new_slide()
std_title(s, "마네킹 — 살인마 스킬 4종")
skills = [
    ("더미 설치",
     ["설치 가능 구역 리스트를 순회",
      "가장 가깝고 유효 거리 내 지점 자동 선택",
      "전이·교란 스킬의 기반 오브젝트"]),
    ("전이",
     ["더미만 렌더링하는 전용 카메라 별도 구성",
      "맵 안개에 의한 시야 차단 해결",
      "→ 안개 무시 커스텀 쉐이더 제작·적용"]),
    ("교란",
     ["전이 애니메이션을 생존자가 파악하는 문제",
      "교란 사용 시 모든 더미에 동시 애니메이션",
      "→ 실제 전이와 구분 불가 → 심리전 확장"]),
    ("위장",
     ["1인칭 시점의 색적 능력 보완",
      "3인칭 시점 전환 기능 추가",
      "인근 거리 내 살인마 존재 시 생존자 효음 차단"]),
]
sx = 0.8
for title, bullets in skills:
    rect(s, sx, 3.0, 7.8, 13.0, CARD)
    rect(s, sx, 3.0, 7.8, 0.07, RGBColor(0xFF, 0x6B, 0x6B))
    box(s, title, sx+0.2, 3.2, 7.2, 0.9, size=15, bold=True, color=RGBColor(0xFF, 0x6B, 0x6B))
    by = 4.3
    for b in bullets:
        box(s, "•  " + b, sx+0.2, by, 7.2, 0.9, size=11, color=WHITE)
        by += 0.9
    sx += 8.2

# 14. 마네킹 — 스폰 최적화
s = new_slide()
std_title(s, "마네킹 — 스폰 시스템 최적화", "네트워크 트래픽 50% 감소")
metric_pair(s, "전송 데이터", "6개 값", "3개 값", 1.5, 3.5, 9)
card(s, "분석",
     ["기존: position(x,y,z) + rotation(x,y,z) = 6개 전송",
      "살인마 position.y = 108f 고정",
      "생존자 position.y = 1f 고정",
      "rotation.x, rotation.z 항상 0 고정",
      "→ 실제로 변하는 값: x, z, rotation.y 3개뿐"],
     11.0, 3.0, 21.5, 7.5)
code = (
    "public (Vector3 pos, Quaternion rot)\n"
    "    GetKillerSpawnData()\n"
    "{\n"
    "    Vector3 d = killerPoints[SelectedIndex];\n"
    "    return (\n"
    "        new Vector3(d.x, 108f, d.z),\n"
    "        Quaternion.Euler(0, d.y, 0));\n"
    "}"
)
box(s, code, 11.0, 11.0, 21.5, 7.0, size=13, font=CODE, color=ACCENT, bg=CARD)

# 15. 마네킹 — 인터페이스 vs 추상 클래스
s = new_slide()
std_title(s, "마네킹 — 인터페이스 vs 추상 클래스", "상호작용 시스템 & 증강 시스템")

rect(s, 0.8,  3.0, 15.8, 14.5, CARD)
rect(s, 0.8,  3.0, 15.8, 0.07, ACCENT)
box(s, "상호작용 시스템 → 인터페이스", 1.1, 3.2, 15.0, 0.9, size=14, bold=True, color=ACCENT)
items_if = [
    "파괴 / 정화 / 뛰어넘기  ← 수평 확장 관계",
    "파괴·정화: NetworkBehaviour 상속 필요",
    "뛰어넘기: MonoBehaviour로 충분",
    "C# 단일 상속 제약 → 추상 클래스로 묶기 불가",
    "→ 인터페이스 선택",
]
for i, t in enumerate(items_if):
    box(s, ("•  " if not t.startswith("→") else "   ") + t,
        1.1, 4.3 + i * 1.0, 15.0, 0.9, size=12,
        color=YELLOW if t.startswith("→") else WHITE)

code_if = (
    "public interface IInteractable\n"
    "{\n"
    "    string GetInteractionText(\n"
    "               GameObject interactor);\n"
    "    string GetAnimationTrigger();\n"
    "    void   OnInteract(\n"
    "               GameObject gameObject);\n"
    "}"
)
box(s, code_if, 1.1, 9.7, 15.0, 7.0, size=12, font=CODE, color=ACCENT, bg=RGBColor(0x20,0x20,0x34))

rect(s, 17.2, 3.0, 15.8, 14.5, CARD)
rect(s, 17.2, 3.0, 15.8, 0.07, RGBColor(0xFF,0x6B,0x6B))
box(s, "증강 시스템 → 추상 클래스", 17.5, 3.2, 15.0, 0.9, size=14, bold=True, color=RGBColor(0xFF,0x6B,0x6B))
items_abs = [
    "NetworkBehaviour 반드시 상속 (호스트·클라이언트 동기화)",
    "기능이 계층적으로 확장 ← 수직 확장 관계",
    "→ 추상 클래스 선택",
    "Activate(): 미구현 시 치명적 → 추상 함수로 강제화",
    "Awake():  방어 코드 (선택적) → 가상 함수로 오버라이드 허용",
]
for i, t in enumerate(items_abs):
    box(s, ("•  " if not t.startswith("→") else "   ") + t,
        17.5, 4.3 + i * 1.0, 15.0, 0.9, size=12,
        color=YELLOW if t.startswith("→") else WHITE)

code_abs = (
    "public abstract class AbilityBase\n"
    "    : NetworkBehaviour\n"
    "{\n"
    "    // 미구현 시 치명적 → 추상 함수\n"
    "    public abstract void Activate();\n\n"
    "    // 방어 코드 → 가상 함수\n"
    "    protected virtual void Awake()\n"
    "    { this.enabled = false; }\n"
    "}"
)
box(s, code_abs, 17.5, 9.7, 15.0, 7.0, size=12, font=CODE, color=ACCENT, bg=RGBColor(0x20,0x20,0x34))

# 16. 마네킹 — Occlusion Culling + Jitter
s = new_slide()
std_title(s, "마네킹 — 렌더링 최적화 & 네트워크 안정화")
metric_pair(s, "렌더링 성능", "13 fps", "63.6 fps", 1.0, 3.0, 10)
card(s, "Occlusion Culling 선택 이유",
     ["AI 생성 3D 메시 다수 → vertex 복잡",
      "맵 규모 상당 → Frustum Culling 한계",
      "카메라 뷰 내 가려진 오브젝트가 많은 구조",
      "5주 일정 → 개별 LOD보다 공간 단위 Occlusion Culling이",
      "비용 대비 효과 높다고 판단"],
     12.0, 3.0, 21.0, 7.0)
card(s, "Jitter 완화",
     ["발생: 특정 증강 획득 시 Jitter 발생",
      "원인: 서버·클라이언트가 서로 다른 시간축으로",
      "       상태 변경 → 네트워크 지연 누적",
      "해결: RPC로 상태값 추가 확정 (방어 코드)",
      "       동기화 기준을 서버로 일원화",
      "       증강 적용 시점을 동기화 확인 이후로 지연"],
     12.0, 10.5, 21.0, 7.0)

# ══════════════════════════════════════════════════════════════════════════════
# 프라가라흐 섹션
# ══════════════════════════════════════════════════════════════════════════════
section_slide("프라가라흐",
              "2D 웹 액션 RPG  |  Phaser · Django · Vue.js  |  2인  |  3개월",
              RGBColor(0xAA, 0x7F, 0xFF))

# 17. 프라가라흐 — 역할
s = new_slide()
std_title(s, "프라가라흐 — 담당 역할")
role_table(s, [
    ["보스 패턴 설계",  "본인", "Min-Heap 기반 패턴 관리 + Dictionary 리팩토링 (OCP)"],
    ["스폰 시스템",     "본인", "안전 구역 기반 설계 — 매 프레임 충돌 연산 제거"],
    ["컷씬 시스템",     "본인", "비트 연산 기반 시청 기록 관리 — DB 컬럼 고정, O(1) 조회"],
    ["Scene 로딩",      "본인", "지연 로딩 전환 — 유령 오브젝트 및 중복 import 제거"],
    ["Django REST API", "본인", "SQLite3 기반 서버 구축 — 동시성 한계 분석 후 적합성 판단"],
    ["코드 구조화",     "본인", "AI 생성 비정형 코드 분석 및 버그 역행 추적 해결"],
])

# 18. 프라가라흐 — Min-Heap 보스 패턴
s = new_slide()
std_title(s, "프라가라흐 — Min-Heap 기반 보스 패턴", "균형 있는 패턴 배치 + Dictionary 리팩토링")
options_boss = [
    ("랜덤 함수", "강한 패턴이 연속 → 억빠/억까 발생", "✗", RED),
    ("쿨타임 기반", "특정 시점에 강한 패턴 집중 사이클", "✗", RED),
    ("Min-Heap + 우선순위", "패턴이 균형 있게 혼합 / O(log N) 삽입", "✓", GREEN),
]
hx, hy = 0.8, 3.0
for label, desc, adopt, ac in options_boss:
    rbg = CARD if adopt == "✗" else RGBColor(0x20, 0x38, 0x28)
    rect(s, hx, hy, 31.5, 1.4, rbg)
    box(s, adopt, hx+0.2,  hy+0.3, 1.2, 0.8, size=16, color=ac, bold=True, align=PP_ALIGN.CENTER)
    box(s, label, hx+1.8,  hy+0.1, 9.0, 0.8, size=13, color=ac, bold=True)
    box(s, desc,  hx+11.2, hy+0.1, 20.5, 0.8, size=12, color=WHITE if adopt == "✗" else GREEN)
    hy += 1.5
card(s, "Dictionary 리팩토링 (OCP)",
     ["Before: switch문 — 새 패턴 추가 시 switch 수정 필요",
      "After: BossPatterns[bossName][priority] 구조",
      "새 패턴 추가 = 딕셔너리에 항목 하나 추가",
      "→ 기존 코드 변경 없음"],
     0.8, 7.5, 15.5, 6.5)
code = (
    "const BossPatterns = {\n"
    "  coffin: {\n"
    "    999: { key:'thunder',   cool:2.5, init:1  },\n"
    "    1:   { key:'summons',   cool:90,  init:5  },\n"
    "    2:   { key:'fireshoot', cool:30,  init:10 },\n"
    "  },\n"
    "  vampire: {\n"
    "    999: { key:'batswarm',  cool:2.5, init:1  },\n"
    "    1:   { key:'summons',   cool:65,  init:23 },\n"
    "  },\n"
    "};"
)
box(s, code, 17.0, 7.5, 16.0, 10.0, size=12, font=CODE, color=ACCENT, bg=CARD)

# 19. 프라가라흐 — 안전 구역 스폰 + 비트 연산
s = new_slide()
std_title(s, "프라가라흐 — 안전 구역 스폰 & 컷씬 비트 연산")
card(s, "안전 구역 기반 스폰",
     ["문제: Phaser에서 충돌 오브젝트가 겹쳐 스폰되면 분리 불가",
      "기존 시도: 매 프레임 밀어내기 → 느림 + 연쇄 충돌",
      "해결: 미리 설정한 안전 좌표 리스트에서 스폰 위치 선택",
      "      겹침 발생 시 양쪽을 반대 방향으로 동일 거리 이동",
      "효과: 매 프레임 충돌 연산 제거 / 빠른 위치 재설정"],
     0.8, 3.0, 15.5, 8.0)
card(s, "컷씬 시청 기록 — 비트 연산",
     ["문제: 씬마다 개별 컬럼 → 씬 추가 시 DB 스키마 변경 필요",
      "해결: cutScene 정수 필드 1개 / 씬 번호 = 비트 위치",
      "효과: DB 컬럼 수 고정 / 조회 O(1) / 스키마 변경 불필요"],
     0.8, 11.5, 15.5, 5.5)
code_spawn = (
    "// 안전 구역 스폰\n"
    "// 1. 미리 설정한 safeZones 리스트\n"
    "// 2. 겹침 감지 시 반대 방향으로 이동\n"
)
code_bit = (
    "// 컷씬 비트 플래그\n"
    "if ((cutScene & (1 << sceneNum)) == 0)\n"
    "{\n"
    "    cutscene.play(introScript);\n"
    "    cutScene += 1 << sceneNum;\n"
    "}"
)
box(s, code_spawn, 17.0, 3.0,  16.0, 5.5, size=13, font=CODE, color=ACCENT, bg=CARD)
box(s, code_bit,   17.0, 9.0,  16.0, 7.0, size=13, font=CODE, color=ACCENT, bg=CARD)

# 20. 프라가라흐 — Scene 로딩 + API 서버
s = new_slide()
std_title(s, "프라가라흐 — Scene 로딩 최적화 & API 서버")
card(s, "Phaser Scene 로딩 최적화",
     ["문제: Phaser 게임 객체 생성 시 모든 씬 동시 로드",
      "→ 씬별 오브젝트가 유령 상태로 잔존 / import key 중복 워닝",
      "해결: 게임 객체 생성 시 로드 씬만 포함",
      "      각 씬에서 필요한 데이터를 지연 로드",
      "효과: 유령 오브젝트 제거 / 중복 워닝 해소"],
     0.8, 3.0, 15.5, 8.0)
card(s, "Django REST API — SQLite3 적합성 판단",
     ["SQLite3: write-lock 기반 동시성 한계 인지",
      "바로 구조 변경 대신 먼저 실제 쓰기 빈도 분석",
      "→ 1인용 게임, 대부분 요청이 조회",
      "→ 쓰기: 회원가입 + 저장 시점만 / 최대 동시 접속 ~1,000명",
      "판단: SQLite3로 충분히 대응 가능"],
     0.8, 11.5, 15.5, 5.5)
card(s, "구조 판단 흐름",
     ["① 문제 인지: SQLite3 동시성 한계",
      "② 실제 부하 분석: 쓰기 빈도 측정",
      "③ 규모 확인: 배포·테스트 대상 SSAFY 서울 캠퍼스 한정",
      "④ 결론: 현재 규모에서는 SQLite3 적합",
      "→ 과도한 기술 전환 없이 근거 있는 판단"],
     17.0, 3.0, 16.0, 14.0)

# 21. 프라가라흐 — AI 코드 구조화
s = new_slide()
std_title(s, "프라가라흐 — AI 생성 코드 구조화 & 버그 해결")
card(s, "배경",
     ["팀원이 AI 코딩 도구로 대량의 코드 생성",
      "컨텍스트·컨벤션 미정의 → 높은 coupling",
      "AI에 다시 의존하면 더 큰 문제 발생 자명",
      "→ 프로젝트 전체 구조 직접 분석"],
     0.8, 3.0, 15.5, 5.5)
card(s, "CASE 1 — 스킬 시스템",
     ["증상: 스킬 포인트 투자 시 특정 스킬만 미적용",
      "원인: 일부 스킬이 씬에서 직접 바인딩",
      "      추상 클래스 상속 시 base 누락",
      "해결: 각 스킬 동작 흐름도 작성",
      "      이상적 구조 + 주석 포함 코드 문서화",
      "      → 팀원에게 리팩토링 방향 제시"],
     0.8, 9.0, 15.5, 8.5)
card(s, "CASE 2 — 데이터 동기화",
     ["증상: 맵 이동·재로그인 시 해금 스킬 초기화",
      "       UI에서는 미해금 표시인데 스킬 사용 가능",
      "방법: 코드량 방대 → 전체 로그 분석 불가",
      "      관련 함수·변수를 역방향 추적(역행 분석)",
      "원인: Vue→Phaser 로드 시 더미 계정 레벨 값이",
      "      API 데이터로 덮어씌워져야 하는데 누락",
      "해결: 해당 필드 수정 + 의심 부분 목록 팀 공유"],
     17.0, 3.0, 16.0, 14.5)

# ══════════════════════════════════════════════════════════════════════════════
# 22. 종합 성과
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
std_title(s, "종합 성과", "3개 프로젝트 핵심 수치 & 역량")

metrics = [
    ("렌더링\n(마네킹)", "13 fps", "63.6 fps"),
    ("드로우콜\n(헤이터)", "3,500",  "180"),
    ("메모리\n(헤이터)",  "60 mb",  "15 mb"),
    ("트래픽\n(마네킹)",  "기준치", "50% 감소"),
]
mx = 0.8
for label, bef, aft in metrics:
    metric_pair(s, label, bef, aft, mx, 3.0, 7.5)
    mx += 8.0

keywords = [
    "설계 판단 근거 언어화",
    "OCP · CRTP · Min-Heap 자료구조 실전 적용",
    "Photon 멀티플레이 · Jitter 완화 경험",
    "WebGL · Occlusion Culling 렌더링 최적화",
    "AI 생성 코드 구조화 및 역행 분석 디버깅",
    "기획·개발 양면 경험 — 기획 총괄 후 직접 구현",
]
box(s, "핵심 역량", 0.8, 12.5, 32.0, 0.8, size=14, bold=True, color=ACCENT)
kx = 0.8
for i, kw in enumerate(keywords):
    rect(s, kx, 13.5, 10.5, 1.1, CARD)
    box(s, kw, kx+0.3, 13.65, 9.9, 0.85, size=11, color=WHITE)
    kx += 10.9
    if i == 2:
        kx = 0.8

prs.save(r"C:\Users\SSAFY\Desktop\Portfolio\포트폴리오개선안\프로젝트_포트폴리오.pptx")
print("저장 완료: 프로젝트_포트폴리오.pptx")
print(f"총 슬라이드 수: {len(prs.slides)}")
